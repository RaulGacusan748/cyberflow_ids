from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(r"C:\Users\user\cyberflow_ids")
REPORTS = ROOT / "reports"
PRESENTATIONS = ROOT / "presentations"
PORTAL = ROOT / "submissions" / "portal_ready"

PRESENTATIONS.mkdir(parents=True, exist_ok=True)
PORTAL.mkdir(parents=True, exist_ok=True)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(22)


def add_image_slide(prs: Presentation, title: str, image_path: Path, caption: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.6), Inches(1.2), width=Inches(11.8), height=Inches(5.4))
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(6.7), Inches(12.0), Inches(0.5))
    tf = tx.text_frame
    tf.text = caption
    tf.paragraphs[0].font.size = Pt(14)


def build_technical_deck(out_path: Path) -> None:
    prs = Presentation()
    add_title_slide(
        prs,
        "CyberFlow IDS Technical Presentation",
        "Dr. Raul C. Gacusan | AIM Capstone | July 2026",
    )
    add_bullet_slide(prs, "Problem and Objective", [
        "Detect malicious network traffic with low false positives and low latency",
        "Frame: binary classification for intrusion detection",
        "Goal: reproducible training-to-deployment IDS workflow",
    ])
    add_bullet_slide(prs, "Dataset and Data Quality", [
        "Dataset: CIC-IDS-2017 derivative traffic",
        "Processed sample: 126,038 rows, 54 columns",
        "Quality checks: null/infinity cleansing, duplicate scan, outlier profiling",
    ])
    add_bullet_slide(prs, "Preprocessing Pipeline", [
        "Label remap: benign/normal -> 0, attack -> 1",
        "Metadata exclusion and numeric-only feature projection",
        "StandardScaler persisted for production parity",
    ])
    add_bullet_slide(prs, "Feature Engineering and Selection", [
        "Leakage shield removes near-proxy features (|r| > 0.98)",
        "Zero-variance feature filtering applied",
        "Final model feature space: 51 numeric features",
    ])
    add_bullet_slide(prs, "Model Tournament", [
        "Candidates: Decision Tree, Random Forest, XGBoost, Naive Bayes",
        "Metrics: Accuracy, Precision, Recall, F1, Train Time",
        "Champion selected by strongest F1 performance",
    ])
    add_image_slide(
        prs,
        "Model Comparison",
        REPORTS / "model_comparison.png",
        "Tournament metrics visualization from training run.",
    )
    add_bullet_slide(prs, "Champion Performance", [
        "XGBoost won tournament (F1 = 0.9868, Accuracy = 0.9956)",
        "Artifacts saved: models/intrusion_detector.joblib, models/scaler.joblib",
        "Inference validated through dashboard and replay simulation",
    ])
    add_image_slide(
        prs,
        "SHAP Explainability",
        REPORTS / "shap_summary_top20.png",
        "Top feature contributions driving intrusion predictions.",
    )
    add_bullet_slide(prs, "Deployment and Reproducibility", [
        "Local FastAPI app implemented for serving predictions",
        "GitHub public repo with release tag: v1.0-submission",
        "Complete submission bundle with reports, code, and evidence",
    ])
    prs.save(str(out_path))


def build_business_deck(out_path: Path) -> None:
    prs = Presentation()
    add_title_slide(
        prs,
        "CyberFlow IDS Executive Business Presentation",
        "AI-Driven Risk Reduction for Network Security Operations",
    )
    add_bullet_slide(prs, "Executive Problem Statement", [
        "Security teams face high alert volume and analyst fatigue",
        "Missed attacks and false positives both create business risk",
        "Need: faster, more reliable triage support",
    ])
    add_bullet_slide(prs, "Solution Overview", [
        "CyberFlow IDS classifies network flow as benign or attack",
        "Production-ready ML pipeline with explainability outputs",
        "Supports SOC triage decisions in near-real time",
    ])
    add_bullet_slide(prs, "Business Value", [
        "Reduce false-positive investigation overhead",
        "Improve analyst productivity and response prioritization",
        "Support faster containment of high-risk events",
    ])
    add_bullet_slide(prs, "Performance Snapshot", [
        "Champion model: XGBoost",
        "Accuracy: 99.56% | F1: 98.68%",
        "Low-latency prediction observed in live replay tests",
    ])
    add_bullet_slide(prs, "Risk and Governance", [
        "Bias and fairness analysis completed in final report",
        "Leakage controls and explainability artifacts are in place",
        "Human-in-the-loop action policy recommended",
    ])
    add_image_slide(
        prs,
        "Explainability for Trust",
        REPORTS / "shap_bar_top20.png",
        "SHAP bar summary used to explain model behavior to stakeholders.",
    )
    add_bullet_slide(prs, "Implementation Readiness", [
        "Artifacts versioned and reproducible via GitHub tag",
        "FastAPI deployment implemented for local serving",
        "Documentation package prepared for audit and handover",
    ])
    add_bullet_slide(prs, "Roadmap", [
        "Phase 1: controlled SOC pilot and monitoring",
        "Phase 2: drift checks, periodic retraining, threshold tuning",
        "Phase 3: expanded telemetry integration and policy automation",
    ])
    add_bullet_slide(prs, "Investment and ROI Narrative", [
        "Lower triage time translates to reduced operational cost",
        "Better detection quality reduces incident impact risk",
        "Reusable pipeline minimizes future model update costs",
    ])
    add_bullet_slide(prs, "Decision and Next Steps", [
        "Approve pilot deployment in controlled environment",
        "Track KPI deltas: false positives, triage time, response speed",
        "Scale after stable governance and performance checks",
    ])
    prs.save(str(out_path))


def main() -> None:
    tech_path = PRESENTATIONS / "CyberFlow_IDS_Technical_Presentation_Rebuilt.pptx"
    biz_path = PRESENTATIONS / "CyberFlow_IDS_Business_Presentation_Rebuilt.pptx"

    build_technical_deck(tech_path)
    build_business_deck(biz_path)

    portal_tech = PORTAL / "Dr_Raul_C_Gacusan_CyberFlow_IDS_Assignment_Technical_Presentation.pptx"
    portal_biz = PORTAL / "Dr_Raul_C_Gacusan_CyberFlow_IDS_Assignment_Business_Presentation.pptx"

    portal_tech.write_bytes(tech_path.read_bytes())
    portal_biz.write_bytes(biz_path.read_bytes())

    print("Generated:")
    print(tech_path)
    print(biz_path)
    print(portal_tech)
    print(portal_biz)


if __name__ == "__main__":
    main()
